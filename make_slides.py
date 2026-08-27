"""Build an editable PowerPoint deck about the eddy-reproduction collaboration.

A high-level, narrative deck for showing colleagues an example of a Claude Code
request and how it went. Re-runnable: edit the content below and re-run.

  /Users/severinf/Applications/anaconda3/bin/python make_slides.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "out"
ASSETS = ROOT / "slides" / "assets"
DECK = ROOT / "slides" / "eddy_collaboration_deck.pptx"

# ---- palette ----
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)   # brick red
BLUE = RGBColor(0x1E, 0x3A, 0xE0)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)   # 16:9


def new_deck():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def band(slide, color=NAVY, height=Inches(1.15)):
    """Top color band for a header."""
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def set_run(p, text, size, color=NAVY, bold=False, italic=False, font="Calibri"):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return r


def header(slide, title, color=NAVY):
    band(slide, color)
    tf = textbox(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.85),
                 anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_run(p, title, 30, WHITE, bold=True)


def bullets(slide, items, l, t, w, h, size=18, gap=8):
    tf = textbox(slide, l, t, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        # support (lead, rest) tuples for bold lead-in
        if isinstance(it, tuple):
            lead, rest = it
            set_run(p, lead, size, ACCENT, bold=True)
            set_run(p, rest, size, GREY)
        else:
            set_run(p, "•  ", size, ACCENT, bold=True)
            set_run(p, it, size, GREY)
    return tf


def img_fit(slide, path, box_l, box_t, box_w, box_h, align="center", valign="middle"):
    """Place an image scaled to fit inside a box, preserving aspect ratio."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    bw, bh = box_w, box_h
    if bw / bh > ar:      # box wider -> fit height
        h = bh; w = int(bh * ar)
    else:
        w = bw; h = int(bw / ar)
    if align == "center":
        l = box_l + (bw - w) // 2
    elif align == "left":
        l = box_l
    else:
        l = box_l + (bw - w)
    if valign == "middle":
        t = box_t + (bh - h) // 2
    elif valign == "top":
        t = box_t
    else:
        t = box_t + (bh - h)
    return slide.shapes.add_picture(str(path), l, t, w, h)


def caption(slide, text, l, t, w, color=NAVY, size=14, bold=True, align=PP_ALIGN.CENTER):
    tf = textbox(slide, l, t, w, Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = align
    set_run(p, text, size, color, bold=bold)


# ============================================================ slides
def build():
    prs = new_deck()

    # ---- 1. Title ----
    s = blank(prs)
    from pptx.enum.shapes import MSO_SHAPE
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    bg.shadow.inherit = False
    tf = textbox(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(2.2))
    p = tf.paragraphs[0]
    set_run(p, "Reproducing a published ocean-eddy result\nwith Claude Code",
            40, WHITE, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(18)
    set_run(p2, "An example of a request we can make — and how it went", 22,
            RGBColor(0xBF, 0xD3, 0xE6), italic=True)
    p3 = tf.add_paragraph(); p3.space_before = Pt(30)
    set_run(p3, "Séverin Fournier   ·   HR-SSS / 2026 RFI", 18,
            RGBColor(0xBF, 0xD3, 0xE6))

    # ---- 2. The request (verbatim) ----
    s = blank(prs)
    header(s, "The request — in my own words")
    tf = textbox(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(2.6))
    p = tf.paragraphs[0]
    set_run(p, "“The overall goal for me is to make an distribution histogram "
               "from this paper’s figure 2a and b. So the plan should be for you "
               "to download all the data used for this figure (in situ salinity "
               "data, only surface), process the data, remake the maps 2a and 2b to "
               "check and then make the histograms”", 24, NAVY, italic=True)
    tf2 = textbox(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(1.5))
    p = tf2.paragraphs[0]
    set_run(p, "Paper:  ", 16, GREY, bold=True)
    set_run(p, "Mo et al. (2024), JGR: Oceans — “A Global Assessment of "
               "Eddy-Induced Salinity Anomalies and Salt Transport by Eddy "
               "Movement”  (10.1029/2023JC020382)", 16, GREY)
    p2 = tf2.add_paragraph(); p2.space_before = Pt(10)
    set_run(p2, "That one paragraph was the whole brief. Everything below followed "
                "from it.", 16, ACCENT, italic=True)

    # ---- 3. The plan ----
    s = blank(prs)
    header(s, "The plan Claude proposed")
    intro = textbox(s, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.7))
    set_run(intro.paragraphs[0],
            "A 5-step pipeline — fully scripted and re-runnable:", 18, NAVY, bold=True)
    steps = [
        ("1  Download  ", "the source data: ~2M+ in-situ salinity profiles (NOAA "
         "World Ocean Database) + the satellite eddy atlas (META3.1exp)."),
        ("2  Extract & QC  ", "keep each profile’s surface (top 10 m) salinity, "
         "applying the paper’s quality-control rules."),
        ("3  Collocate  ", "match every profile to the nearest same-day eddy; tag it "
         "inside-eddy (d < R) or background (d > 2R)."),
        ("4  Grid  ", "on a 2°×2° grid, anomaly = mean(inside-eddy) − "
         "mean(background), separately for anticyclonic (AE) and cyclonic (CE)."),
        ("5  Reproduce  ", "remake maps 2a/2b to check against the paper, then build "
         "the distribution histogram."),
    ]
    bullets(s, steps, Inches(0.9), Inches(2.15), Inches(11.6), Inches(4.8),
            size=19, gap=14)

    # ---- 4. Issues we hit ----
    s = blank(prs)
    header(s, "Issues we bumped into — and how we solved them", ACCENT)
    issues = [
        ("Paper was paywalled  ", "→ you downloaded the PDF and dropped it in "
         "the folder."),
        ("Eddy data not where expected  ", "→ it lives on AVISO+, not "
         "Copernicus; its FTP was blocked, so we found the public AVISO gateway."),
        ("Download throttled to ~10 KB/s  ", "→ the server capped single "
         "streams; switched to a chunked parallel downloader (~10× faster)."),
        ("Data-format bug  ", "→ the WOD files store each variable with its own "
         "layout; first parser mis-aligned depth & salinity. Found and fixed, then "
         "validated on real files."),
        ("Unphysical outliers  ", "→ a few sparse grid cells blew up; added a "
         "min-profiles-per-cell rule and a sanity range on salinity values."),
    ]
    bullets(s, issues, Inches(0.9), Inches(1.5), Inches(11.7), Inches(5.4),
            size=18, gap=13)
    note = textbox(s, Inches(0.9), Inches(6.75), Inches(11.7), Inches(0.5))
    set_run(note.paragraphs[0],
            "Claude diagnosed and fixed each of these; most needed no input from me.",
            15, NAVY, italic=True)

    # ---- 5. My inputs ----
    s = blank(prs)
    header(s, "What I steered — every input I gave")
    intro = textbox(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.6))
    set_run(intro.paragraphs[0],
            "Claude drove; I made a handful of decisions at key forks:", 18,
            NAVY, bold=True)
    inputs = [
        ("Data access  ", "provided the paper PDF and confirmed the eddy-data route "
         "(Copernicus login, then AVISO gateway)."),
        ("Scope  ", "full global, 1993–2019, surface (top 10 m) only."),
        ("Quality threshold  ", "reviewed the maps; chose to keep the stricter "
         "min-10-profiles-per-cell (cleaner, closer to the paper)."),
        ("Equator band  ", "asked to blank ±5° latitude, as the paper does."),
        ("Mission-concept framing  ", "defined the detection floor 0.2/√7 pss "
         "(0.2 pss single-pass accuracy over a 7-day eddy decorrelation)."),
        ("Figure styling  ", "colours (red+blue→purple), labels, mid-latitude "
         "band, and units (psu→pss)."),
    ]
    bullets(s, inputs, Inches(0.9), Inches(2.0), Inches(11.7), Inches(5.2),
            size=18, gap=12)

    # ---- 6. Maps side by side ----
    s = blank(prs)
    header(s, "Result — reproduced maps vs. the paper")
    colw = Inches(6.1)
    caption(s, "Paper — Mo et al. (2024) Fig 2a/2b", Inches(0.5), Inches(1.35),
            colw, color=GREY)
    img_fit(s, ASSETS / "paper_fig2_ab.png", Inches(0.5), Inches(1.8),
            colw, Inches(2.4), valign="top")
    caption(s, "Our reproduction", Inches(6.75), Inches(1.35), colw, color=ACCENT)
    img_fit(s, OUT / "fig2ab_check.png", Inches(6.75), Inches(1.75),
            colw, Inches(5.3), valign="top")
    note = textbox(s, Inches(0.5), Inches(4.5), colw, Inches(2.3))
    for i, it in enumerate([
        ("Same signature.  ", "Anticyclonic eddies (AE) are saltier across the "
         "mid-latitudes; cyclonic eddies (CE) are the mirror image."),
        ("Independently built  ", "from raw public data — not traced from the "
         "paper."),
    ]):
        p = note.paragraphs[0] if i == 0 else note.add_paragraph()
        p.space_after = Pt(10)
        lead, rest = it
        set_run(p, lead, 16, ACCENT, bold=True); set_run(p, rest, 16, GREY)

    # ---- 7. Histograms ----
    s = blank(prs)
    header(s, "Result — the distribution histogram (the deliverable)")
    img_fit(s, OUT / "fig2_histogram.png", Inches(0.4), Inches(1.5),
            Inches(6.3), Inches(4.6), valign="top")
    img_fit(s, OUT / "fig2_histogram_combined.png", Inches(6.9), Inches(1.5),
            Inches(6.3), Inches(4.6), valign="top")
    note = textbox(s, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.9))
    p = note.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p, "AE anomalies skew positive (saltier), CE negative (fresher) — "
               "the paper’s headline result, recovered from scratch. Right: "
               "both pooled, with the 0.2/√7 pss detection floor.", 15, NAVY)

    # ---- 8. Takeaways ----
    s = blank(prs)
    header(s, "Takeaways")
    cards = [
        ("Scale handled for me", "2.96M salinity profiles processed; ~90 GB "
         "downloaded, parsed, and gridded — all scripted."),
        ("Reproducible", "Every step is a numbered, re-runnable script; the whole "
         "result regenerates from one command."),
        ("I stayed in control", "Claude drove the work and flagged decisions; I "
         "steered at ~6 key forks."),
        ("From one paragraph", "A single plain-language request → a verified "
         "reproduction of a published figure + a new mission-concept analysis."),
    ]
    from pptx.enum.shapes import MSO_SHAPE
    x0, y0 = Inches(0.7), Inches(1.5)
    cw, ch = Inches(5.85), Inches(1.9)
    gap = Inches(0.25)
    for i, (title, body) in enumerate(cards):
        r, c = divmod(i, 2)
        cx = x0 + c * (cw + gap)
        cy = y0 + r * (ch + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch)
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = NAVY; card.line.width = Pt(1)
        card.shadow.inherit = False
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        p = tf.paragraphs[0]
        set_run(p, title, 20, ACCENT, bold=True)
        p2 = tf.add_paragraph(); p2.space_before = Pt(8)
        set_run(p2, body, 16, GREY)

    # time strip along the bottom
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7),
                               Inches(5.85), Inches(11.98), Inches(1.35))
    strip.fill.solid(); strip.fill.fore_color.rgb = NAVY
    strip.line.fill.background(); strip.shadow.inherit = False
    tf = strip.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]
    set_run(p, "Time  —  ", 17, RGBColor(0xBF, 0xD3, 0xE6), bold=True)
    set_run(p, "My hands-on effort: just a handful of decisions and figure "
               "reviews.   ", 16, WHITE)
    p2 = tf.add_paragraph(); p2.space_before = Pt(6)
    set_run(p2, "Downloading was the long pole (unattended). Parallel / chunked "
                "downloads made it fast — the ~14 GB eddy atlas took minutes; the "
                "~90 GB of salinity data would have too, had we parallelized it "
                "from the start (lesson learned).", 16, WHITE)
    p3 = tf.add_paragraph(); p3.space_before = Pt(6)
    set_run(p3, "Claude’s active work (processing + all figure iterations, "
                "excluding downloads): under ~30 minutes.", 16,
            RGBColor(0xFF, 0xD9, 0x9E), bold=True)

    DECK.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DECK))
    print("wrote", DECK, "with", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    build()
